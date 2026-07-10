#!/usr/bin/env python3
"""
Deck de planificación del Plan Manuel Belgrano (v2).

Fuente: planilla de Carmen Caro Solís "Proyecto Cáñamo / Proyecto Belgrano"
(Google Sheets, exportada a .xlsx). Genera un .pptx SIN diseño (solo textos
editables — el diseño lo hace Guille), 16:9.

Uso:
    python3 build_deck.py <planilla.xlsx> <salida.pptx>

Hojas usadas:
    - "Plan Belgrano"  → misión + 5 programas (proyectos, owners, indicadores)
    - "Planificación"  → modelo de madurez APLICADO (5 dimensiones × 6 etapas,
                          con estado real: logrado / PENDIENTE)
    - "V01"            → plantilla madre (de acá salen los criterios de Gate)
    - "Impactos"       → catálogo de métricas
"""

import re
import sys

import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt

ETAPAS = [
    ("Etapa 0", "Formulación"),
    ("Etapa 1", "Viabilidad"),
    ("Etapa 2", "Piloto"),
    ("Etapa 3", "Demostración"),
    ("Etapa 4", "Escalamiento"),
    ("Etapa 5", "Plataforma de Desarrollo Territorial"),
]
DIMENSIONES = ["Técnica", "Regulatoria", "Comercial", "Financiera", "Institucional"]


# ---------------- parsing ----------------

def parse_plan_belgrano(ws):
    """Devuelve (mision, [ {titulo, objetivo, proyectos:[(nombre, owner, indicadores)]} ])."""
    mision = (ws["B2"].value or "").strip()
    programas = []
    actual = None
    for row in ws.iter_rows(min_row=3, max_row=60):
        c = row[2].value  # col C
        d = row[3].value  # col D
        if c and str(c).strip().startswith("Programa"):
            actual = {"titulo": str(c).strip(), "objetivo": "", "proyectos": []}
            programas.append(actual)
            continue
        if c and actual and not actual["objetivo"]:
            obj = str(c).strip()
            actual["objetivo"] = re.sub(r"^Objetivo:\s*", "", obj)
            continue
        if d and actual:
            owner = str(row[4].value).strip() if row[4].value else ""
            inds = [str(row[i].value).strip() for i in (5, 6, 7) if row[i].value]
            actual["proyectos"].append((str(d).strip(), owner, inds))
    return mision, programas


def parse_celda_dimension(texto):
    """Una celda de 'Planificación' → (pregunta, [(titulo_kr, pendiente:bool)])."""
    texto = texto.strip()
    partes = texto.split("•")
    pregunta = partes[0].strip()
    krs = []
    for b in partes[1:]:
        b = b.strip()
        titulo = re.split(r"\s+—\s+", b, maxsplit=1)[0]
        titulo = re.sub(r"\s+", " ", titulo).strip().rstrip(".")
        krs.append((titulo, "PENDIENTE" in b))
    return pregunta, krs


def parse_planificacion(ws):
    """Devuelve matriz[dimension][etapa_idx] = (pregunta, krs)."""
    matriz = {}
    for r in range(4, 9):
        dim = (ws.cell(row=r, column=1).value or "").strip()
        if dim not in DIMENSIONES:
            continue
        celdas = []
        for c in range(2, 8):  # B..G = etapas 0..5
            v = ws.cell(row=r, column=c).value
            celdas.append(parse_celda_dimension(str(v)) if v else ("", []))
        matriz[dim] = celdas
    return matriz


def parse_gates(ws_v01):
    """Gates de la plantilla madre V01.
    Devuelve gates[gate_idx 0..4][dimension] = criterio (texto corto).
    Columnas C,E,G,I,K = gates 1..5; filas 4..7 = Técnica..Financiera."""
    gates = [dict() for _ in range(5)]
    filas = {4: "Técnica", 5: "Regulatoria", 6: "Comercial", 7: "Financiera"}
    for r, dim in filas.items():
        for gi, col in enumerate([3, 5, 7, 9, 11]):
            v = ws_v01.cell(row=r, column=col).value
            if v:
                txt = re.sub(r"\s+", " ", str(v)).strip()
                txt = re.split(r"\s*-->\s*|\s*→\s*", txt)[0].strip()
                gates[gi][dim] = txt
    return gates


def parse_impactos(ws):
    cols = {"Ambiental": [], "Social": [], "Económico": []}
    claves = list(cols)
    for row in ws.iter_rows(min_row=2, max_row=10):
        for i, k in enumerate(claves):
            v = row[i].value
            if v:
                cols[k].append(str(v).strip().rstrip("."))
    return cols


# ---------------- pptx ----------------

BLANK = 6  # layout en blanco


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def slide(self, titulo, sub=None):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[BLANK])
        tb = s.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = titulo
        p.font.size = Pt(30)
        p.font.bold = True
        if sub:
            p2 = tf.add_paragraph()
            p2.text = sub
            p2.font.size = Pt(14)
        return s

    def cuerpo(self, s, x=0.6, y=1.5, w=12.1, h=5.7):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        return tf

    @staticmethod
    def parr(tf, texto, size=14, bold=False, nivel=0, antes=6):
        p = tf.paragraphs[0] if (len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs) else tf.add_paragraph()
        p.text = texto
        p.font.size = Pt(size)
        p.font.bold = bold
        p.level = nivel
        p.space_before = Pt(antes)
        return p

    def save(self, path):
        self.prs.save(path)


def main(xlsx_path, out_path):
    wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    mision, programas = parse_plan_belgrano(wb["Plan Belgrano"])
    matriz = parse_planificacion(wb["Planificación"])
    gates = parse_gates(wb["V01"])
    impactos = parse_impactos(wb["Impactos"])

    d = Deck()

    # 1 — Portada
    s = d.slide("Plan Manuel Belgrano", "Proyecto Cáñamo — Documento de planificación")
    tf = d.cuerpo(s, y=3.2, h=2.5)
    d.parr(tf, "Plataforma de desarrollo regenerativo a base de cáñamo industrial", 20, bold=True)
    d.parr(tf, "Planificación: Lic. Carmen Caro Solís · Flora Cáñamo Neuquino", 13)
    d.parr(tf, "Julio 2026 — borrador de trabajo (contenido editable, sin diseño)", 12)

    # 2 — Misión
    s = d.slide("La misión")
    tf = d.cuerpo(s, y=2.2)
    d.parr(tf, mision, 22)

    # 3 — Overview de programas
    s = d.slide("Cinco programas, un solo plan")
    tf = d.cuerpo(s)
    for prog in programas:
        d.parr(tf, prog["titulo"], 17, bold=True, antes=12)
        d.parr(tf, prog["objetivo"], 13, nivel=1)

    # 4–8 — Un slide por programa
    for prog in programas:
        s = d.slide(prog["titulo"], prog["objetivo"])
        tf = d.cuerpo(s, y=1.8, h=5.4)
        for nombre, owner, inds in prog["proyectos"]:
            linea = nombre.rstrip(".")
            extras = []
            if owner:
                extras.append(f"Owner/Partner: {owner}")
            if inds:
                extras.append("Indicadores: " + " · ".join(i.rstrip(".") for i in inds))
            d.parr(tf, "— " + linea, 14, bold=True, antes=8)
            if extras:
                d.parr(tf, "   " + "  |  ".join(extras), 11, nivel=1, antes=0)

    # 9 — Intro modelo de madurez
    s = d.slide("Modelo de madurez", "6 etapas × 5 dimensiones, con gates de inversión entre etapas")
    tf = d.cuerpo(s, y=1.9)
    d.parr(tf, "Etapas: " + "  →  ".join(f"{n} {t}" for n, t in ETAPAS), 14, bold=True)
    d.parr(tf, "Dimensiones evaluadas en paralelo en cada etapa: " + " · ".join(DIMENSIONES), 14, antes=14)
    d.parr(tf, "Cada etapa define indicadores verificables (KRs). Un Gate entre etapas responde una sola "
               "pregunta: ¿hay evidencia suficiente para invertir en la etapa siguiente?", 14, antes=14)
    d.parr(tf, "Estado de avance según la planilla vigente: los indicadores ya alcanzados se marcan ✔; "
               "los no alcanzados, ⏳ PENDIENTE.", 13, antes=14)

    # 10–15 — Un slide por etapa
    for ei, (num, nombre) in enumerate(ETAPAS):
        s = d.slide(f"{num} — {nombre}")
        tf = d.cuerpo(s, y=1.5, h=5.8)
        for dim in DIMENSIONES:
            pregunta, krs = matriz.get(dim, [("", [])] * 6)[ei]
            if not pregunta and not krs:
                continue
            d.parr(tf, dim + " — " + pregunta.split("\n")[0], 13, bold=True, antes=8)
            for titulo, pendiente in krs:
                marca = "⏳" if pendiente else "✔"
                d.parr(tf, f"{marca} {titulo}", 11, nivel=1, antes=1)
        if ei < 5 and gates[ei]:
            d.parr(tf, f"GATE {ei + 1} → {ETAPAS[ei + 1][0]} ({ETAPAS[ei + 1][1]}):", 12, bold=True, antes=10)
            for dim, criterio in gates[ei].items():
                d.parr(tf, f"{dim}: {criterio}", 10, nivel=1, antes=1)

    # 16 — Métricas de impacto
    s = d.slide("Métricas de impacto", "Qué se mide en cada dimensión del triple impacto")
    tf = d.cuerpo(s, y=1.9)
    for k, items in impactos.items():
        d.parr(tf, f"Impacto {k}", 16, bold=True, antes=12)
        d.parr(tf, " · ".join(items), 13, nivel=1)

    # 17 — Próximos pasos (derivados de los PENDIENTE de Etapas 0–1)
    s = d.slide("Próximos pasos", "Pendientes clave de las Etapas 0 y 1 según la planilla vigente (borrador editable)")
    tf = d.cuerpo(s, y=1.8)
    pendientes = []
    for dim in DIMENSIONES:
        for ei in (0, 1):
            _, krs = matriz.get(dim, [("", [])] * 6)[ei]
            for titulo, pendiente in krs:
                if pendiente:
                    pendientes.append((dim, titulo))
    for dim, titulo in pendientes:
        d.parr(tf, f"— [{dim}] {titulo}", 12, antes=3)

    d.save(out_path)
    print(f"OK: {out_path} ({len(d.prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main(sys.argv[1], sys.argv[2])
